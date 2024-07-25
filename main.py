import json
import re
import os
from pptx import Presentation

PRESENTATIONS_DIRECTORY = "Presentations"
JSON_DIRECTORY = 'JSONs'


def sanitize_and_format_filename(filename):
    # Remove file extension
    name, ext = os.path.splitext(filename)

    # Remove any destructive characters
    sanitized_name = re.sub(r'[<>:"/\\|?*=+]', '', name).strip()

    # Convert to title case and replace spaces with underscores
    formatted_name = '_'.join(word.capitalize() for word in sanitized_name.split())

    # Return the new filename with the original extension
    return formatted_name + ext


def rename_files_in_directory(directory_path):
    # List all files in the directory
    for filename in os.listdir(directory_path):
        file_path = os.path.join(directory_path, filename)

        new_filename = sanitize_and_format_filename(filename)
        new_file_path = os.path.join(directory_path, new_filename)

        # Rename the file
        os.rename(file_path, new_file_path)
        print(f"Renamed: {filename} -> {new_filename}")

def extract_presentation_text(pptx_path, presentation_title):
    presentation = Presentation(pptx_path)
    presentation_data = {
        # "presentation_title": presentation.core_properties.title or "Untitled Presentation",
        "presentation_title": presentation_title,
        "slides": []
    }

    slide_number = 1
    # for slide in presentation.slides:
    #     slide_title = None
    #     slide_texts = []
    #
        # for shape in slide.shapes:
        #     if shape.has_text_frame:    # shape.has_text_frame is True if this shape can contain text
        #
        #         if shape.text:    # If shape's text isn't empty
        #             text = shape.text
        #             text = text.replace("\n", "REAPLACEMENT")
        #             text = text.replace("\v", "REAPLACEMENT")
        #
        #             if shape == slide.shapes[0]:
        #                 slide_title = text  # Assuming the first shape contains the title
        #             else:
        #                 slide_texts.append(text)

    for slide in presentation.slides:
        text_shapes = []

        for shape in slide.shapes:
            if shape.has_text_frame and shape.text:  # shape.has_text_frame is True if this shape can contain text
                text = re.sub(r"[\n\v]+", " ", shape.text)
                text = re.sub("[ ]{2,}", " ", text).strip()
                text_shapes.append(text)

        slide_title = text_shapes[0] if text_shapes else "Untitled Slide" # Assuming that the first text shape contains the title
        slide_texts = text_shapes[1:] if len(text_shapes) > 1 else []

        slide_content = {
            "slide_number": slide_number,
            "slide_title": slide_title or "Untitled Slide",
            "slide_text": "BREAK".join(slide_texts)
        }
        slide_number += 1

        presentation_data["slides"].append(slide_content)

    return presentation_data


def save_to_json(data, json_path):
    with open(json_path, 'w', encoding='utf-8') as json_file:
        json.dump(data, json_file, ensure_ascii=False, indent=4)


def PPTX_to_JSON(directory_path):
    # List all files in the directory
    for filename in os.listdir(directory_path):

        file_path = os.path.join(directory_path, filename)
        file_title = os.path.splitext(filename)[0]

        # Extract presentation text and save to JSON
        extracted_text = extract_presentation_text(file_path, file_title)
        json_path = os.path.join(JSON_DIRECTORY, f"{file_title}.json")
        save_to_json(extracted_text, json_path)

    print(f"Presentations reformatted and saved to {JSON_DIRECTORY}.")


rename_files_in_directory(PRESENTATIONS_DIRECTORY)
PPTX_to_JSON(PRESENTATIONS_DIRECTORY)